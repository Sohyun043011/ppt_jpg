from PyQt5 import uic
import os, shutil
import sys
from PyQt5.QtWidgets import *
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Pt
from PyQt5.QtGui import QPainter, QPen, QColor, QPixmap , QIcon
from PyQt5.QtCore import QTimer, Qt
import urllib.request
import webbrowser
from comtypes import client
import shutil

dataImage_default_path="C:\\Server\\Gachi\\Qname\\dataImage"

form_class = uic.loadUiType("ppt_to_jpg.ui")[0] # ppt_to_jpg.ui(xml 형식)에서 레이아웃 및 텍스트 설정값 조정

class MyWindow(QMainWindow, form_class):
    def __init__(self):
        super().__init__()
        self.setupUi(self)
        self.initSetting()
        self.setFixedSize(1600, 850) # 창 사이즈 고정
        self.setWindowTitle('화상회의실 관리 프로그램') # 프로그램 Title 설정
        self.setWindowIcon(QIcon('./wrench.png')) # 프로그램 아이콘 설정
        
        
        self.connect_count=0 # 연결 시도 횟수 설정
        
        self.nameLinkBtn.setDisabled(True) # 초기 링크 버튼 비활성화
        self.wallLinkBtn.setDisabled(True) # 초기 링크 버튼 비활성화
        
        self.chrome_path="C:/Program Files (x86)/Google/Chrome/Application/chrome.exe %s" # Chrome 설치 위치
        
        self.timer = QTimer(self)
        self.timer.setInterval(6000)
        self.timer.timeout.connect(self.update_network)
        self.timer.start()
        
        self.LayoutTab.setCurrentIndex(0)       #layout1 기본 선택 되게 설정
       
        
        
        
    def initSetting(self):
        # 서식 이미지 파일
        pixmap1 = QPixmap("./양식1.jpg")
        pixmap2 = QPixmap("./양식2.jpg")
        pixmap3 = QPixmap("./양식3.jpg")

        self.image1.setPixmap(QPixmap(pixmap1))
        self.image2.setPixmap(QPixmap(pixmap2))
        self.image3.setPixmap(QPixmap(pixmap3))
        
        self.statusLabel.setText('연결 없음')
        self.statusLabel.setStyleSheet("color : red")
        
        # 서식 버튼 설정
        self.radioBtn_1.toggled.connect(self.onClicked)
        self.radioBtn_2.toggled.connect(self.onClicked)
        self.radioBtn_3.toggled.connect(self.onClicked)
        self.radioBtn_1.setChecked(True)
        self.selectForm.clicked.connect(self.onClickSelect)
        self.selectForm_2.clicked.connect(self.onClickSelect)
        
        self.createBtn.clicked.connect(self.createBtn_clicked)
        self.createBtn_2.clicked.connect(self.createBtn_clicked)
        self.deleteBtn.clicked.connect(self.deleteBtn_clicked)
        self.deleteBtn_2.clicked.connect(self.deleteBtn_clicked)
        
        # 활성화 버튼 설정
        self.wallLinkBtn.clicked.connect(self.onWallOpenClick)
        self.wallActivBtn.clicked.connect(self.onWallActivClick)
        self.nameLinkBtn.clicked.connect(self.onNameOpenClick)
        self.nameActivBtn.clicked.connect(self.onNameActivClick)
        
        
        self.set_style()
    
    def set_style(self):
        with open("update_style", 'r') as f:
            self.setStyleSheet(f.read())    
        
    #  라디오버튼을 통해 서식을 결정한 경우, 해당 서식 파일의 경로를 설정
    def onClicked(self):
        global pptx_fpath
        radioBtn = self.sender()
        if radioBtn.text() == '서식1':
            pptx_fpath = os.path.dirname(os.path.abspath('양식1.pptx'))+'\\양식1.pptx'
        elif radioBtn.text() == '서식2':
            pptx_fpath = os.path.dirname(os.path.abspath('양식2.pptx'))+'\\양식2.pptx'
        elif radioBtn.text() == '서식3':
            pptx_fpath =  os.path.dirname(os.path.abspath('양식3.pptx'))+'\\양식3.pptx'
            
    #  회의실 책상 나타낸 도형 그리기 함수
    def paintEvent(self,event):
        qp = QPainter()
        qp.begin(self)
        #그리기 함수의 호출부분
        self.drawRectangles(qp)
        qp.end()
        
    # 회의실 책상 나타낸 도형 그리기 함수
    def drawRectangles(self,qp):
        qp.setBrush(QColor(255, 136, 79))
        qp.setPen(QPen(QColor(255, 136, 79), 3))
        qp.drawRect(529, 265, 221, 30)
        qp.drawRect(529, 265, 30, 361)
        qp.drawRect(720, 265, 30, 361)
        qp.drawRect(270, 230, 20, 191)
        
    # 사용자가 선택한 양식을 가져와서 ppt 생성하는 함수
    def makePPT(self,directory,subject,pptx_fpath,inputValue):
        prs = Presentation(pptx_fpath)                  # 양식 선택시 불러옴
        #슬라이드 1~15 까지 돌면서,기재된 값(이름,직위)을 text_on_shape() 에 넘겨줌.
        for i in range(15):
            slide = prs.slides[i]
            inputVal = inputValue[i+1]
            shapes = slide.shapes
            self.text_on_shape(shapes,inputVal)
        prs.save(directory+'\\'+subject+'.pptx')            # /팀명/subject이름/ 폴더에 만든 ppt 저장.
    
    # makePPT로 만든 ppt 를 JPG로 변환해주는 함수 .1~15개 슬라이드를 1~30 장으로 변환한다. 
    def makeJPG(self,directory,subject):
        ppt = client.CreateObject('Powerpoint.Application')
        # print(directory+"\\"+subject+".pptx")
        subject=subject+".pptx"
        ppt.Presentations.Open(os.path.join(directory,subject))
        ppt.ActivePresentation.Export(directory, 'JPG')
        ppt.ActivePresentation.Close()
        ppt.Quit()
        
        # 슬라이드 1~15를 돌면서, 각각 1->16, 2->17,... 로 복사해준다.JPG파일 이름 또한 `슬라이드1.jpg` 를 `1.jpg` 로 변경함.
        for i in range(15):
            os.rename(directory+"\\슬라이드" + str(i+1) + ".JPG", directory + "\\" + str((i+1)) + '.jpg')
            shutil.copyfile(directory+"\\"+str((i+1)) + '.jpg',directory+"\\"+str(15+(i+1)) + '.jpg')
    
    # 사용자가 [이름, 직위] 입력 박스에 적은 값을 받아오는 함수  
    def inputValue(self):
        inputValue= {}
        # layout1 인 경우(이름1개, 직위1개)
        if currentIndex==0:
            for i in range(1,15):
                nameChild = self.findChild(QLineEdit,"InputName_%d" % (i)).text()
                namePos = self.findChild(QLineEdit,"InputPos_%d" % (i)).text()
                inputValue[i]=[nameChild,namePos]
            inputValue[15]=['','']    #(15번째 슬라이드는 항상 공란이므로, 공백으로 저장함.)
        else:
            # layout2 인 경우(이름2개, 직위2개)
            for i in range(1,15):
                nameChild_L = self.findChild(QLineEdit,"InputName_%d_L" % (i)).text()
                nameChild_R = self.findChild(QLineEdit,"InputName_%d_R" % (i)).text()
                namePos_L = self.findChild(QLineEdit,"InputPos_%d_L" % (i)).text()
                namePos_R = self.findChild(QLineEdit,"InputPos_%d_R" % (i)).text()
                inputValue[i] = [[nameChild_L,nameChild_R],[namePos_L,namePos_R]]
            inputValue[15]=[['',''],['','']]
        return inputValue
    
    # [서식 찾기] 버튼 누를 때 실행되는 함수.파일 탐색기를 열어 사용자가 선택한 양식의 path를 가져옴.
    def onClickSelect(self):
        global pptx_fpath
        currentIndex = self.LayoutTab.currentIndex()
        select_file = QFileDialog.getOpenFileName(self)             # 파일탐색기 열림
        pptx_fpath = select_file[0]                                 # 해당 파일의 경로(ex) C://어쩌고~//양식7.pptx
        if pptx_fpath=='':
            QMessageBox.about(self,"message","파일이 선택되지 않았습니다.다시 선택해주세요.")
        else:
            self.findFormLabel.setText("선택 서식 : \n"+pptx_fpath) if currentIndex==0 else  self.findFormLabel2.setText("선택 서식 : \n"+pptx_fpath)
            
    # 사용자가 입력한 값(이름,직위)을 새로 생성한 ppt에 적용하는 함수
    def text_on_shape(self,shapes,inputVal):
        # shapes : 한 슬라이드 안의 구조들
        for shape in shapes:
            # layout1 인 경우(name,pos)
            if shape.name=="name" or shape.name=="pos":
                # shape 내 텍스트 프레임 선택하기 & 기존 값 삭제하기
                text_frame = shape.text_frame
                p = text_frame.paragraphs[0]
                font_size = p.runs[0].font.size.pt      # 선택한 양식 텍스트상자의 폰트 사이즈
                font_color = p.runs[0].font.color       # 선택한 양식 텍스트상자의 폰트 컬러
                font_bold = p.runs[0].font.bold         # 선택한 양식 텍스트상자의 폰트 굵기
                font_name = p.runs[0].font.name         # 선택한 양식 텍스트상자의 폰트 이름
                text_frame.clear()                      # 폰트 설정 다 받아온 뒤, 텍스트 상자 비워줌
                # 정렬 설정 : 중간정렬
                p.alighnment = PP_ALIGN.CENTER   
                run = p.add_run()
                run.text = inputVal[0] if shape.name=="name" else inputVal[1]       #이름이면 inputVal 배열의 첫번째 인자를 넣고, pos 이면 배열의 두번째 인자를 넣음
                # 위에서 저장해 둔 폰트 설정을 그대로 적용함.
                font = run.font 
                font.name = font_name
                font.size = Pt(font_size)
                if font_bold==True:
                    # bold 설정 되어있다면
                    font.bold = font_bold
            
                if font_color.type!=None:
                    # 블랙아닌 경우 - SCHEME,RGB
                    # 1. type==SCHEME인 경우
                    if font_color.theme_color!=0:
                        font.color.theme_color=font_color.theme_color
                        font.color.brightness = font_color.brightness           #폰트 색상의 투명도가 있는 경우, 투명도 설정
                        
                    else:
                        # 2. type==RGB인 경우
                        # font_color.rgb 는 FF0000 처럼 출력 되는데, 이를 255,0,0 처럼 변환함. RGBColor(255,0,0)
                        font.color.rgb = RGBColor(int(str(font_color.rgb)[0:2],16),int(str(font_color.rgb)[2:4],16),int(str(font_color.rgb)[4:6],16))
                        font.color.brightness = font_color.brightness
                else:
                    # 블랙인 경우,type이 None임.
                    font.color.rgb = RGBColor(0,0,0)
                    
            # layout2 인 경우(name1,name2,pos1,pos2)
            elif shape.name=='name1' or shape.name=='name2' or shape.name=='pos1' or shape.name=='pos2':
                # shape 내 텍스트 프레임 선택하기 & 기존 값 삭제하기
                text_frame = shape.text_frame
                p = text_frame.paragraphs[0]
                font_size = p.runs[0].font.size.pt
                font_color = p.runs[0].font.color
                font_bold = p.runs[0].font.bold
                font_name = p.runs[0].font.name
                text_frame.clear()
                # 정렬 설정 : 중간정렬
                p.alighnment = PP_ALIGN.CENTER   
                run = p.add_run()
                if shape.name=="name1":
                    run.text = inputVal[0][0]
                elif shape.name=="name2":
                    run.text = inputVal[0][1]
                elif shape.name=="pos1":
                    run.text = inputVal[1][0]
                elif shape.name=="pos2":
                    run.text = inputVal[1][1]
                    
                font = run.font 
                font.name = font_name
                font.size = Pt(font_size)
                if font_bold==True:
                    # bold 설정 되어있다면
                    font.bold = font_bold
                if font_color.type!=None:
                    # 블랙아닌 경우 - SCHEME,RGB
                    # 1. type==SCHEME인 경우
                    if font_color.theme_color!=0:
                        font.color.theme_color=font_color.theme_color
                        font.color.brightness = font_color.brightness
                    else:
                        # 2. type==RGB인 경우
                        font.color.rgb = RGBColor(int(str(font_color.rgb)[0:2],16),int(str(font_color.rgb)[2:4],16),int(str(font_color.rgb)[4:6],16))
                        font.color.brightness = font_color.brightness
                else:
                    # 블랙인 경우,
                    font.color.rgb = RGBColor(0,0,0)
                
    def disableBtn(func):
        def wrapper(self):
            self.nameLinkBtn.setDisabled(True) # 초기 링크 버튼 비활성화
            self.wallLinkBtn.setDisabled(True) # 초기 링크 버튼 비활성화
            func(self)
        return wrapper            
    
    # create 버튼 클릭시 이벤트       
    def createBtn_clicked(self):
        # 현재 tab이 어디인지 확인 (layout1 : 0 or layout2 : 1)
        global currentIndex
        currentIndex = self.LayoutTab.currentIndex()
        
        # /팀이름/subject/ 로 폴더 생성
        subject = self.subject.text() if currentIndex==0 else self.subject_2.text()         # 폴더 이름
        
        if subject=='':
            # 공백인 경우 alert 띄움
            QMessageBox.about(self,"message","폴더명을 입력해주세요.")
        else:
            deptLabel = self.deptName.currentText() if currentIndex==0 else self.deptName_2.currentText()   # 부서명
            # directory = os.path.join(dataImage_default_path,deptLabel,subject)     # 디렉토리 경로
            directory = os.getcwd()+"\\"+deptLabel+"\\"+subject
            inputValue = self.inputValue()      #사용자가 입력한 정보
            if not os.path.exists(directory):
                os.makedirs(directory)
                # 폴더 생성 후, ppt 생성
                self.makePPT(directory,subject,pptx_fpath,inputValue)       #디렉토리 경로,폴더이름,선택한 양식경로,입력값
                self.makeJPG(directory,subject)                             #디렉토리 경로, 폴더이름
            else: 
                # 이미 있는 폴더인 경우, alert 띄움
                QMessageBox.about(self,"message",subject+"는 이미 있는 폴더입니다. 다른 이름을 설정해주세요.")
                
            # create 버튼실행이 완료된 경우, subject 입력박스와 선택 서식 표출 박스 비워줌.
            self.subject.clear()
            self.subject_2.clear()
            self.findFormLabel.clear()
        
    
    def deleteBtn_clicked(self):
        # create 버튼 클릭시 이벤트, new_dialog를 통해 새 창에서 지울 수 있는 UI 구성
        self.dialog=QDialog(self)
        self.dialog.setWindowModality(Qt.ApplicationModal)
        self.dialog.setWindowTitle('삭제')
        self.dialog.setFixedSize(430,370)
        self.dialog.selectLabel=QLabel('부서 폴더를 선택해주세요.', self.dialog)
        self.dialog.selectedDept=''
        self.dialog.selectLabel.move(10,20)
        self.dialog.combo_box=QComboBox(self.dialog)
        
        self.dialog.listWidget=QListWidget(self.dialog)
        self.dialog.listWidget.setGeometry(10,80,400,200) 
        
        self.dialog.deleteBtn=QPushButton('삭제하기',self.dialog)
        self.dialog.deleteBtn.setGeometry(10,300,100,50) 
        self.dialog.deleteBtn.clicked.connect(self.underFolderDelete)
        
        #combobox list 추가 (dataImage에 존재하는 모든 디렉토리명 가져오기)
        
        folder_list = os.listdir(dataImage_default_path)
        for folder in folder_list:
            self.dialog.combo_box.addItem(folder)
        self.dialog.combo_box.move(10,40)
        self.dialog.combo_box.activated[str].connect(self.onActived)
        
        #combobox list 선택 시 하위 폴더들 리스트로 가져오기
        self.dialog.show()
    
    def onActived(self, text): # 부서 선택 시 하위 폴더들 표출
        self.dialog.listWidget.clear()
        
        under_folder_list = os.listdir(os.path.join(dataImage_default_path, text))
        for folder in under_folder_list:
            QListWidgetItem(folder, self.dialog.listWidget)

    def underFolderDelete(self):
        temp_path=[item.text() for item in self.dialog.listWidget.selectedItems()]
        if len(temp_path)==0:
            return

        folderPath=os.path.join(dataImage_default_path, self.dialog.combo_box.currentText() ,temp_path[0])

        
        if os.path.exists(folderPath): # 폴더가 존재할 때
            # 폴더 상위폴더가 dataImage일 경우(부서 카테고리를 지우려고 하는 경우)
            if folderPath.split('\\')[-2]=='dataImage':
                QMessageBox.warning(self, '알림', '부서 카테고리는 삭제할 수 없습니다. 다시 시도해주세요.')
                return
            if '/'.join(folderPath.split('\\')[:-2])!='C:/Server/Gachi/Qname/dataImage':
                QMessageBox.warning(self, '알림', '다른 디렉토리에 존재하는 파일은 삭제할 수 없습니다. 다시 시도해주세요.')
                return
            
            buttonReply=QMessageBox.information(self,'알림',f'정말로 {temp_path[0]} 폴더를 삭제하시겠습니까?', QMessageBox.Yes|QMessageBox.No, QMessageBox.No)
            if buttonReply==QMessageBox.Yes:
                try:
                    shutil.rmtree(folderPath) # 폴더 하위에 파일의 유무에 관계없이 무조건 삭제
                    QMessageBox.information(self,'알림','폴더와 하위 파일들이 삭제되었습니다.')
                    lst_modelIndex=self.dialog.listWidget.selectedIndexes()
                    for modelIndex in lst_modelIndex:
                        self.dialog.listWidget.model().removeRow(modelIndex.row())
                except Exception as e:
                    QMessageBox.warning(self, '알림', '다른 프로세스가 파일을 사용 중이기 때문에 해당 폴더를 제거할 수 없습니다.')
        else:
            QMessageBox.information(self,'알림','폴더가 존재하지 않습니다.')
      
    def ping(self, ip):
        if self.connect_count>=4: #시도 횟수 2번 이상이면 0번으로 갱신 후 stop (한번 시도마다 2개의 ip에 대해 조사)
            self.timer.stop()
            self.connect_count=0
            return False
        self.connect_count+=1
        try:
            print('다음으로 연결 중: http://'+ip)
            urllib.request.urlopen('http://'+ip, timeout=1)
            return True
        except urllib.request.URLError as err:
            return False
    
    def update_network(self): # 스마트월 ip와 디지털명패 ip 각각의 연결성을 확인 후 상태 표시
        wall_ip="192.168.0.60" #스마트월 ip
        name_ip="192.168.0.103/Qname/empMain.aspx?readImage=ok" #디지털명패 ip
        self.nameLinkBtn.setDisabled(True)
        self.wallLinkBtn.setDisabled(True)
        
        if self.ping(wall_ip):
            self.statusLabel.setText('스마트월 연결 성공')
            self.statusLabel.setStyleSheet("color : blue;")
            self.wallLinkBtn.setDisabled(False)
            self.timer.stop()
            self.connect_count=0
        elif self.ping(name_ip):
            self.statusLabel.setText('디지털명패 연결 성공')
            self.statusLabel.setStyleSheet("color : blue;")
            self.nameLinkBtn.setDisabled(False)
            self.timer.stop()
            self.connect_count=0
        else:
            self.statusLabel.setText('연결 없음')
            self.statusLabel.setStyleSheet("color : red;")
            
    @disableBtn        
    def onNameActivClick(self): # 디지털명패 활성화 버튼 눌렀을 때 onclick function
        QMessageBox.information(self,'알림','디지털명패가 활성화되었습니다. 아래 링크 버튼이 활성화가 될 때까지 잠시만 기다려주세요.')
        os.startfile('enable.bat.lnk')
        self.timer.start()

    @disableBtn
    def onWallActivClick(self): # 스마트월 활성화 버튼 눌렀을 때 onclick function
        os.startfile('disable.bat.lnk')
        QMessageBox.information(self,'알림','스마트월이 활성화되었습니다. 아래 링크 버튼이 활성화가 될 때까지 잠시만 기다려주세요.')
        self.timer.start()
    
    def onWallOpenClick(self):
        webbrowser.get(self.chrome_path).open("192.168.0.60")
    
    def onNameOpenClick(self):
        webbrowser.get(self.chrome_path).open("http://192.168.0.103/Qname/empMain.aspx?readImage=ok")
    
    
    
if __name__ == "__main__":
    try:
        app = QApplication(sys.argv)
        myWindow = MyWindow()
        myWindow.show()
        app.exec_()
    except Exception as e:
        print(e)