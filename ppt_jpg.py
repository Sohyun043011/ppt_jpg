import os
import sys
import time
from PyQt5.QtWidgets import QApplication, QWidget, QPushButton, QTextEdit, QLabel, QCheckBox, QButtonGroup, QRadioButton
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Pt
from comtypes import client
from PyQt5.QtWidgets import QWidget, QApplication
from PyQt5.QtGui import QPainter, QPen, QColor, QBrush, QPixmap
from PyQt5.QtCore import Qt
import socket
import subprocess
import threading

class BtnCreatePpt(QWidget):
    
    def __init__(self):
        super().__init__()
        self.initUI()

    def initUI(self):
        
        self.btn = QPushButton('Create JPG Files', self)
        self.btn.setGeometry(40, 50,180,30)
        self.btn.clicked.connect(self.btnfunc1)

        self.label_meeting = QLabel('Subject', self)
        self.label_meeting.setGeometry(30,100, 100, 30)
        self.txt_meeting = QTextEdit(self)
        self.txt_meeting.setGeometry(80,100, 150, 30)

        self.label_pos = [0,0,0,0,0,0,0,0,0,0,0,0,0]
        self.label_name = [0,0,0,0,0,0,0,0,0,0,0,0,0]
        self.txt_pos = [0,0,0,0,0,0,0,0,0,0,0,0,0]
        self.txt_name = [0,0,0,0,0,0,0,0,0,0,0,0,0]

        #1~13번 자리 포지션 라벨
        self.label_pos_geometry = [[110,700, 100, 30],[110,600, 100, 30],[110,500, 100, 30],[110,400, 100, 30],[110,300, 100, 30],
        [110,200, 100, 30],[310,100, 100, 30],[660,200, 100, 30],[660,300, 100, 30],[660,400, 100, 30],[660,500, 100, 30],[660,600, 100, 30],[660,700, 100, 30]]
        #1~13번 자리 이름 라벨
        self.label_name_geometry = [[100,740, 100, 30],[100,640, 100, 30],[100,540, 100, 30],[100,440, 100, 30],[100,340, 100, 30],
        [100,240, 100, 30],[300,140, 100, 30],[660,240, 100, 30],[660,340, 100, 30],[660,440, 100, 30],[660,540, 100, 30],[660,640, 100, 30],[660,740, 100, 30]]
        #1~13번 자리 포지션 빈칸
        self.txt_pos_geometry = [[150,700, 100, 30],[150,600, 100, 30],[150,500, 100, 30],[150,400, 100, 30],[150,300, 100, 30],
        [150,200, 100, 30],[350,100, 100, 30],[550,200, 100, 30],[550,300, 100, 30],[550,400, 100, 30],[550,500, 100, 30],[550,600, 100, 30],[550,700, 100, 30]]
        #1~13번 자리 포지션 빈칸
        self.txt_name_geometry = [[150,740, 100, 30],[150,640, 100, 30],[150,540, 100, 30],[150,440, 100, 30],[150,340, 100, 30],
        [150,240, 100, 30],[350,140, 100, 30],[550,240, 100, 30],[550,340, 100, 30],[550,440, 100, 30],[550,540, 100, 30],[550,640, 100, 30],[550,740, 100, 30]]


        for i in range(13):
            self.label_pos[i] = QLabel('Pos' + str(i+1), self)
            self.label_pos[i].setGeometry(self.label_pos_geometry[i][0],self.label_pos_geometry[i][1], self.label_pos_geometry[i][2], self.label_pos_geometry[i][3])
            self.label_name[i] = QLabel('Name' + str(i+1), self)
            self.label_name[i].setGeometry(self.label_name_geometry[i][0],self.label_name_geometry[i][1], self.label_name_geometry[i][2], self.label_name_geometry[i][3])
            self.txt_pos[i] = QTextEdit(self)
            self.txt_pos[i].setGeometry(self.txt_pos_geometry[i][0],self.txt_pos_geometry[i][1], self.txt_pos_geometry[i][2], self.txt_pos_geometry[i][3])
            self.txt_name[i] = QTextEdit(self)
            self.txt_name[i].setGeometry(self.txt_name_geometry[i][0],self.txt_name_geometry[i][1], self.txt_name_geometry[i][2], self.txt_name_geometry[i][3])

        # self.cb1 = QCheckBox('서식1', self)
        # self.cb1.setGeometry(600, 50, 100, 30)
        # self.cb2 = QCheckBox('서식2', self)
        # self.cb2.setGeometry(720, 50, 100, 30)
        # self.cb3 = QCheckBox('서식3', self)
        # self.cb3.setGeometry(840, 50, 100, 30)

        self.rb1 = QRadioButton('서식1', self)
        self.rb1.setGeometry(600, 50, 100, 30)
        self.rb2 = QRadioButton('서식2', self)
        self.rb2.setGeometry(720, 50, 100, 30)
        self.rb3 = QRadioButton('서식3', self)
        self.rb3.setGeometry(840, 50, 100, 30)
        
        # self.cb1.stateChanged.connect(self.check_select)
        # self.cb2.stateChanged.connect(self.check_select)
        # self.cb3.stateChanged.connect(self.check_select)
        # self.cb1.setChecked(True)

        self.rb1.toggled.connect(self.onClicked)
        self.rb2.toggled.connect(self.onClicked)
        self.rb3.toggled.connect(self.onClicked)
        self.rb1.setChecked(True)

        # self.buttongroup = QButtonGroup(self)
        # self.buttongroup.setExclusive(True)
        # self.buttongroup.addButton(self.cb1, 1)
        # self.buttongroup.addButton(self.cb2, 2)
        # self.buttongroup.addButton(self.cb3, 3)
        # self.buttongroup.buttonClicked.connect(self.check_buttongroup)

        self.label_img1 = QLabel('123', self)
        self.label_img1.setGeometry(570, 80, 100, 70)
        self.label_img2 = QLabel('234', self)
        self.label_img2.setGeometry(690, 80, 100, 70)
        self.label_img3 = QLabel('345', self)
        self.label_img3.setGeometry(810, 80, 100, 70)

        pixmap1 = QPixmap("./양식1.jpg")
        pixmap2 = QPixmap("./양식2.jpg")
        pixmap3 = QPixmap("./양식3.jpg")

        self.label_img1.setPixmap(QPixmap(pixmap1))
        self.label_img2.setPixmap(QPixmap(pixmap2))
        self.label_img3.setPixmap(QPixmap(pixmap3))

        self.label_cp = QLabel('모드전환 : 시작', self)
        self.label_cp.setGeometry(300, 850, 190, 50)

        # 디지털 회의실 셋팅
        # wall_ip="192.168.0.7" #비디오월 ip
        # name_ip="192.168.0.103" #스마트명패 ip

        wall_ip="192.168.10.32" #비디오월 ip - LAN 모드
        name_ip="192.168.0.132" #스마트명패 ip - WIFI 모드


        ipaddress=socket.gethostbyname(socket.gethostname())

        cp_btn_name = ''
        if ipaddress == wall_ip:
            cp_btn_name = "비디오월 → 스마트명패"
            self.label_cp.setText("현재모드 : 비디오월")
        elif ipaddress == name_ip:
            cp_btn_name = "스마트명패 → 비디오월"
            self.label_cp.setText("현재모드 : 스마트명패")

        self.cp_btn = QPushButton(cp_btn_name, self)
        self.cp_btn.setGeometry(300, 900, 190, 50)
        self.cp_btn.clicked.connect(lambda :self.btnfunc2(self.cp_btn, self.label_cp))

        self.setGeometry(300, 50, 1200, 900)
        self.setWindowTitle("Smart_Nameplate")
        self.show()

    def onClicked(self):
        global pptx_fpath
        global ex_flag
        radioBtn = self.sender()
        if radioBtn.text() == '서식1':
            pptx_fpath = './양식1.pptx'
            ex_flag = '양식1'
            print('양식1 select')
        elif radioBtn.text() == '서식2':
            pptx_fpath = './양식2.pptx'
            ex_flag = '양식2'
            print('양식2 select')
        elif radioBtn.text() == '서식3':
            pptx_fpath = './양식3.pptx'
            ex_flag = '양식3'
            print('양식3 select')

    def check_buttongroup(self):
        print('call buttongroup function')
    
    def check_select(self, state):
        global pptx_fpath
        global ex_flag
        if state == 1:
            pptx_fpath = os.path.dirname(os.path.abspath('양식1.pptx'))+'\\양식1.pptx'
            ex_flag = '양식1'
            print('양식1 select')
        elif state == 2:
            pptx_fpath = os.path.dirname(os.path.abspath('양식2.pptx'))+'\\양식2.pptx'
            ex_flag = '양식2'
            print('양식2 select')
        elif state == 3:
            pptx_fpath = os.path.dirname(os.path.abspath('양식3.pptx'))+'\\양식3.pptx'
            ex_flag = '양식3'
            print('양식3 select')

    def paintEvent(self, e):
        qp = QPainter()
        qp.begin(self)
        self.draw_rect(qp)
        qp.end()

    def draw_rect(self, qp):
        qp.setBrush(QColor(255, 136, 79))
        qp.setPen(QPen(QColor(255, 136, 79), 3))
        qp.drawRect(280, 200, 50, 570)

        qp.setBrush(QColor(255, 136, 79))
        qp.setPen(QPen(QColor(255, 136, 79), 3))
        qp.drawRect(470, 200, 50, 570)

        qp.setBrush(QColor(255, 136, 79))
        qp.setPen(QPen(QColor(255, 136, 79), 3))
        qp.drawRect(330, 200, 140, 50)

        qp.setBrush(QColor(87, 103, 247))
        qp.setPen(QPen(QColor(87, 103, 247), 3))
        qp.drawRect(300, 800, 190, 50)


    def btnfunc1(self):

        def text_on_shape(shape, input_text, ex_flag, shape_flag, bold = True):

            # shape 내 텍스트 프레임 선택하기 & 기존 값 삭제하기
            text_frame = shape.text_frame
            text_frame.clear()

            # 문단 선택하기
            p = text_frame.paragraphs[0]

            # 정렬 설정 : 중간정렬
            p.alighnment = PP_ALIGN.CENTER   

            # 텍스트 입력 / 폰트 지정
            run = p.add_run()
            run.text = input_text
            font = run.font
            if ex_flag == '양식1':
                if shape_flag == 'name':
                    font.size = Pt(168)
                    font.color.rgb = RGBColor(255,255,255)
                elif shape_flag == 'pos':
                    font.size = Pt(50)
                    font.color.rgb = RGBColor(255,255,255)
            elif ex_flag == '양식2':
                if shape_flag == 'name':
                    font.size = Pt(133)
                    font.color.rgb = RGBColor(0,0,0)
                elif shape_flag == 'pos':
                    font.size = Pt(28)
                    font.color.rgb = RGBColor(255,255,255)
            elif ex_flag == '양식3':
                if shape_flag == 'name':
                    font.size = Pt(133)
                    font.color.rgb = RGBColor(0,0,0)
                elif shape_flag == 'pos':
                    font.size = Pt(28)
                    font.color.rgb = RGBColor(255,255,255)
            font.bold = bold
            font.name = None  # 지정하지 않으면 기본 글씨체로  #  'Arial'

        # os.remove('result.pptx')
        # prs = Presentation()      

        # pptx 파일 열기(양식 선택 포함)

        # pptx_fpath = './양식1.pptx'
        # pptx_fpath = './양식2.pptx'(ppt 파일 얻어오기 - 국내관광팀)

        prs = Presentation(pptx_fpath)


        for index in range(13):

            # 슬라이드 지정하기
            slide_num = index
            slide = prs.slides[slide_num]

            # 슬라이드 내 shape 사전 만들기
            shapes_list = slide.shapes
            shape_index = {}
            
            for i, shape in enumerate(shapes_list):
                shape_index[ shape.name ] = i
            
            print(shape_index)   # {'직사각형(파란부분)': 0, 'logo': 1, 'name': 2, 'pos': 3}

            shape_name = "name"
            shape_select = shapes_list[ shape_index[ shape_name ]]    

            raw_name = self.txt_name[index].toPlainText()
            final_name = ""

            for i in range(len(raw_name)):
                final_name += raw_name[i]
                if i < len(raw_name):
                    final_name += ' '

            text_on_shape(shape_select, final_name, ex_flag, shape_name)

            shape_name = "pos"
            shape_select = shapes_list[ shape_index[ shape_name ]]

            raw_pos = self.txt_pos[index].toPlainText()
            final_pos = ""

            for i in range(len(raw_pos)):
                final_pos += raw_pos[i]
                if i < len(raw_pos):
                    final_pos += ' '

            # text_on_shape(shape_select, self.txt_pos[index].toPlainText(), ex_flag, shape_name)
            text_on_shape(shape_select, final_pos, ex_flag, shape_name)

        prs.save('result.pptx')

        ppt = client.CreateObject('Powerpoint.Application')
        
        # 절대경로 설정
        path = os.path.dirname(os.path.abspath('result.pptx'))
        # print(path+"\\result.pptx") 
        
        #잦은 에러 
        ppt.Presentations.Open(path+"\\result.pptx")
        folder_name = self.txt_meeting.toPlainText()
        ppt.ActivePresentation.Export(path+"\\"+folder_name, 'JPG')
        folder_path = path+"\\"+folder_name
        #foldeR_path : C:\Users\test\Desktop\workspace\ppt_jpg\folder_name
        ppt.ActivePresentation.Close()
        ppt.Quit()

        for i in range(13):
            os.rename(folder_path+"\\슬라이드" + str(i+1) + ".JPG", folder_path + "\\" + str(i+1) + '.jpg')


    def btnfunc2(self, vbutton, vlabel):
        wall_ip="192.168.10.32" #비디오월 ip - LAN 모드
        name_ip="192.168.0.132" #스마트명패 ip - WIFI 모드

        #loop로 모드 전환 확인 필요
        ipaddress=socket.gethostbyname(socket.gethostname())

        if ipaddress == wall_ip:

            # t = threading.Thread(target=swit_func, args=(ipaddress,vbutton, vlabel))
            # t.start()

            print("btnfunc2 Clicked : 비디오월 → 스마트명패 변경 시작")
            os.system("netsh interface set interface \"이더넷 2\" disable")
            os.system("netsh interface set interface \"Wi-Fi\" enable")
            # self.window.after(5000, self.btnfunc2)

            vbutton.setText("스마트명패 → 비디오월") 
            vlabel.setText("현재모드 : 스마트명패")

            # while True:
            #     ipaddress=socket.gethostbyname(socket.gethostname())
            #     if ipaddress == wall_ip:
            #         vlabel.setText("현재모드 : 전환중...")
            #     elif ipaddress == name_ip:
            #         vbutton.setText("스마트명패 → 비디오월") 
            #         vlabel.setText("현재모드 : 스마트명패")
            #         #걍 꺼짐
            #         break

        elif ipaddress == name_ip:
            vlabel.setText("현재모드 : 전환중...")
            print("btnfunc2 Clicked : 스마트명패 → 비디오월 변경 시작")
            os.system("netsh interface set interface \"Wi-Fi\" disable")
            os.system("netsh interface set interface \"이더넷 2\" enable")

            vbutton.setText("비디오월 → 스마트명패")
            vlabel.setText("현재모드 : 비디오월")

            # while True:
            #     ipaddress=socket.gethostbyname(socket.gethostname())
            #     if ipaddress == name_ip:
            #         vlabel.setText("현재모드 : 전환중...")
            #     elif ipaddress == wall_ip:
            #         vbutton.setText("비디오월 → 스마트명패")
            #         vlabel.setText("현재모드 : 비디오월")
            #         #걍 꺼짐
            #         break
        
        def swit_func(self, ipipadress, vvbutton, vvlabel):
            vvbutton.setText("test") 
            time.sleep(1)
            # wall_ip="192.168.10.32" #비디오월 ip - LAN 모드
            # name_ip="192.168.0.132" #스마트명패 ip - WIFI 모드
            # ipipipaddress=socket.gethostbyname(socket.gethostname())
            # if ipipadress == ipipipaddress:
            #     vvbutton.setDisabled(True)
            #     vvlabel.setText("현재모드 : 전환중...")
            # else:
            #     vvbutton.setEnabled(True)
            #     if wall_ip == ipipipaddress:
            #         vvbutton.setText("비디오월 → 스마트명패")
            #         vvlabel.setText("현재모드 : 비디오월")
            #     elif name_ip == ipipipaddress:
            #         vvbutton.setText("스마트명패 → 비디오월") 
            #         vvlabel.setText("현재모드 : 스마트명패")
            #     return
            # time.sleep(1)

if __name__ == '__main__':
    app = QApplication(sys.argv)
    ex = BtnCreatePpt()
    sys.exit(app.exec_())

